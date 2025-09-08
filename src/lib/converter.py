"""ファイル変換処理モジュール

Excel、Word、PDF等の各種ファイルをMarkdown形式に変換する機能を提供します。
v2.3.1で空白行処理機能を追加し、ExcelからMarkdown変換時の空白行を適切に処理します。
"""

import logging
import os
from typing import List, Optional

# Configure module logger
logger = logging.getLogger(__name__)
empty_line_logger = logging.getLogger(f"{__name__}.empty_line")

# Optional imports
try:
    import frontmatter
    _has_frontmatter = True
except ImportError:
    _has_frontmatter = False

# Import empty line handling functions
from .config import get_empty_line_config, EmptyLineConfig


def read_text_file(path: str) -> str:
    """テキストファイルを読み込んで文字列として返す。
    
    Args:
        path: ファイルパス
        
    Returns:
        ファイルの内容
    """
    with open(path, "r", encoding="utf-8") as f:
        return f.read()


def _maybe_markitdown_convert(text: str) -> str:
    """markitdownライブラリが利用可能な場合に変換を試みる。
    
    Args:
        text: 変換対象のテキスト
        
    Returns:
        変換されたテキスト（markitdownが無い場合はそのまま）
    """
    try:
        import markitdown
        md = markitdown.MarkItDown()
        result = md.convert_local(text)
        return result.text_content
    except ImportError:
        return text


# ========================================
# 空白行処理ユーティリティ関数
# ========================================

def is_empty_cell(cell_value) -> bool:
    """セルの値が空白かどうかを判定する。
    
    Args:
        cell_value: セルの値
        
    Returns:
        空白の場合True
    """
    if cell_value is None:
        return True
    if isinstance(cell_value, str) and cell_value.strip() == "":
        return True
    return False


def is_empty_row(row: List) -> bool:
    """行が全て空白セルかどうかを判定する。
    
    Args:
        row: 行データのリスト
        
    Returns:
        全て空白の場合True
    """
    if not row:
        return True
    return all(is_empty_cell(cell) for cell in row)


def is_empty_row_for_table(row_cells: List[str]) -> bool:
    """テーブル用の空白行判定（文字列セルの配列用）。
    
    Args:
        row_cells: テーブル行のセル文字列のリスト
        
    Returns:
        空白行の場合True
    """
    if not row_cells:
        return True
    return all(cell.strip() == "" for cell in row_cells)


def _is_empty_table_line(line: str) -> bool:
    """テーブル行が空白行かどうかを判定する。
    
    Args:
        line: Markdownテーブルの行
        
    Returns:
        空白行の場合True
    """
    if not line or line.strip() == "":
        return True
    
    # Markdownテーブル形式の空白行を判定
    # 例: | | | | (空白セルのみのテーブル行)
    if line.startswith("|") and line.endswith("|"):
        # パイプ区切りの内容を取得
        cells = [cell.strip() for cell in line.split("|")[1:-1]]
        return all(cell == "" for cell in cells)
    
    return False


def filter_consecutive_empty_table_rows(table_lines: List[str]) -> List[str]:
    """テーブル行から連続する空白行を単一の空白行に集約する。
    
    Args:
        table_lines: Markdownテーブルの行のリスト
        
    Returns:
        フィルタリングされた行のリスト
    """
    if not table_lines:
        return table_lines
    
    filtered_lines = []
    prev_was_empty = False
    
    for line in table_lines:
        is_current_empty = _is_empty_table_line(line)
        
        if is_current_empty:
            if not prev_was_empty:
                # 最初の空白行のみ追加
                filtered_lines.append(line)
            # 連続する空白行はスキップ
        else:
            # 非空白行は常に追加
            filtered_lines.append(line)
        
        prev_was_empty = is_current_empty
    
    return filtered_lines


def clean_table_end(table_lines: List[str]) -> List[str]:
    """テーブルの末尾の空白行を削除する。
    
    Args:
        table_lines: Markdownテーブルの行のリスト
        
    Returns:
        末尾空白行が削除された行のリスト
    """
    if not table_lines:
        return table_lines
    
    # 末尾から空白行を削除
    while table_lines and _is_empty_table_line(table_lines[-1]):
        table_lines.pop()
    
    return table_lines


def safe_empty_line_processing(content: str, config: Optional[EmptyLineConfig] = None) -> str:
    """安全な空白行処理を実行する。
    
    Args:
        content: 処理対象のコンテンツ
        config: 空白行処理設定（Noneの場合はデフォルト設定を取得）
        
    Returns:
        処理されたコンテンツ
    """
    try:
        if config is None:
            config = get_empty_line_config()
        
        if not config.enabled:
            return content
        
        # 改行で分割
        lines = content.split('\n')
        
        # 連続する空白行の処理
        if config.remove_consecutive:
            filtered_lines = []
            prev_was_empty = False
            
            for line in lines:
                is_current_empty = line.strip() == ""
                
                if is_current_empty:
                    if config.preserve_single_empty and not prev_was_empty:
                        # 単一空白行は保持
                        filtered_lines.append(line)
                    # 連続する空白行はスキップ
                else:
                    # 非空白行は常に追加
                    filtered_lines.append(line)
                
                prev_was_empty = is_current_empty
            
            lines = filtered_lines
        
        # 末尾の空白行処理
        if config.remove_trailing:
            while lines and lines[-1].strip() == "":
                lines.pop()
        
        return '\n'.join(lines)
        
    except Exception as e:
        # 空白行処理でエラーが発生した場合、元のコンテンツを返す
        empty_line_logger.warning(f"Empty line processing failed, returning original content: {e}")
        return content


def safe_empty_line_table_processing(table_lines: List[str], config: Optional[EmptyLineConfig] = None) -> List[str]:
    """安全なテーブル空白行処理を実行する。
    
    Args:
        table_lines: テーブル行のリスト
        config: 空白行処理設定（Noneの場合はデフォルト設定を取得）
        
    Returns:
        処理されたテーブル行のリスト
    """
    try:
        if config is None:
            config = get_empty_line_config()
        
        if not config.enabled:
            return table_lines
        
        processed_lines = table_lines.copy()
        
        # 連続空白行のフィルタリング
        if config.remove_consecutive:
            processed_lines = filter_consecutive_empty_table_rows(processed_lines)
        
        # 末尾空白行の削除
        if config.remove_trailing:
            processed_lines = clean_table_end(processed_lines)
        
        return processed_lines
        
    except Exception as e:
        # テーブル処理でエラーが発生した場合、元のテーブルを返す
        empty_line_logger.warning(f"Empty line table processing failed, returning original table: {e}")
        return table_lines


# ========================================
# メインの変換関数
# ========================================

def convert_file_to_markdown(path: str) -> str:
    """与えられたファイルを Markdown 文字列に変換して返す。

    Args:
        path: 変換対象ファイルのパス

    Returns:
        Markdown 形式の文字列

    Raises:
        ValueError: 未対応の拡張子の場合
    """
    ext = os.path.splitext(path)[1].lower()
    file_name = os.path.basename(path)

    logger.info(f"ファイル変換開始: {file_name} (形式: {ext})")

    # 空白行処理設定を事前に取得してログ出力
    try:
        empty_line_config = get_empty_line_config()
        if empty_line_config.enabled:
            empty_line_logger.info(f"空白行処理が有効: consecutive={empty_line_config.remove_consecutive}, "
                                 f"trailing={empty_line_config.remove_trailing}, "
                                 f"preserve_single={empty_line_config.preserve_single_empty}")
        else:
            empty_line_logger.info("空白行処理は無効です")
    except Exception as e:
        logger.warning(f"空白行設定の取得に失敗: {e}")

    # --- Markdown ---
    if ext == ".md":
        text = read_text_file(path)
        empty_line_config = get_empty_line_config()

        if _has_frontmatter:
            post = frontmatter.loads(text)
            content = post.content
        else:
            content = text

        # 安全な空白行処理を適用
        content = safe_empty_line_processing(content, empty_line_config)

        logger.info(f"Markdownファイル変換完了: {file_name}")
        return content

    # --- Plain text ---
    if ext == ".txt":
        text = read_text_file(path)
        empty_line_config = get_empty_line_config()

        # 安全な空白行処理を適用
        text = safe_empty_line_processing(text, empty_line_config)

        logger.info(f"テキストファイル変換完了: {file_name}")
        return _maybe_markitdown_convert(text)

    # --- DOCX ---
    if ext == ".docx":
        try:
            import docx  # type: ignore
        except Exception as exc:
            raise RuntimeError("python-docx is required to convert .docx files") from exc

        empty_line_config = get_empty_line_config()
        doc = docx.Document(path)  # type: ignore

        try:
            if empty_line_config.enabled:
                # 空白行処理が有効な場合、段落をフィルタリング
                paragraphs = []
                for p in doc.paragraphs:
                    if p.text and not is_empty_cell(p.text.strip()):
                        paragraphs.append(p.text)
                    elif not empty_line_config.remove_consecutive:
                        # 連続する空白行の除去が無効の場合、空の段落も含める
                        paragraphs.append('')
            else:
                paragraphs = [p.text for p in doc.paragraphs if p.text]

            # 末尾の空白行処理
            if empty_line_config.enabled and empty_line_config.remove_trailing:
                while paragraphs and is_empty_cell(paragraphs[-1]):
                    paragraphs.pop()

            raw = "\n\n".join(paragraphs)
        except Exception as exc:
            # 空白行処理でエラーが発生した場合、元の処理にフォールバック
            logger.warning(f"DOCX空白行処理でエラー、フォールバック実行: {exc}")
            paragraphs = [p.text for p in doc.paragraphs if p.text]
            raw = "\n\n".join(paragraphs)

        logger.info(f"DOCXファイル変換完了: {file_name}")
        return _maybe_markitdown_convert(raw)

    # --- XLSX / XLSM ---
    if ext in [".xlsx", ".xlsm"]:
        try:
            import openpyxl  # type: ignore
        except Exception as exc:
            raise RuntimeError("openpyxl is required to convert .xlsx files") from exc

        # 空白行処理設定を取得
        empty_line_config = get_empty_line_config()

        wb = openpyxl.load_workbook(path, read_only=True, data_only=True)
        sheets_md: List[str] = []

        for sheet in wb.worksheets:
            rows = list(sheet.iter_rows(values_only=True))
            if not rows:
                continue

            # If there are multiple columns, render as a Markdown table using the first row as header
            if any(len(r) > 1 for r in rows if r):
                header = rows[0]
                body = rows[1:]
                # normalize to strings and replace None with empty
                def norm(cell):
                    return "" if cell is None else str(cell)

                header_cells = [norm(c) for c in header]
                # divider
                divider = ["---"] * len(header_cells)
                table_lines = ["| " + " | ".join(header_cells) + " |", "| " + " | ".join(divider) + " |"]

                # 空白行処理の統計情報
                original_row_count = len(body)
                processed_rows = 0
                skipped_empty_rows = 0

                try:
                    for r in body:
                        row_cells = [norm(c) for c in r]
                        # pad/truncate to header length
                        if len(row_cells) < len(header_cells):
                            row_cells += [""] * (len(header_cells) - len(row_cells))

                        # 空白行処理が有効で、空白行の場合はスキップ
                        if empty_line_config.enabled and is_empty_row_for_table(row_cells):
                            skipped_empty_rows += 1
                            continue

                        table_lines.append("| " + " | ".join(row_cells[: len(header_cells)]) + " |")
                        processed_rows += 1

                    # 連続空白行のフィルタリング（空白行処理が有効な場合）
                    if empty_line_config.enabled and empty_line_config.remove_consecutive:
                        table_lines = filter_consecutive_empty_table_rows(table_lines)

                    # 末尾空白行の削除（空白行処理が有効な場合）
                    if empty_line_config.enabled and empty_line_config.remove_trailing:
                        table_lines = clean_table_end(table_lines)

                except Exception as exc:
                    # 空白行処理でエラーが発生した場合、元の処理にフォールバック
                    table_lines = ["| " + " | ".join(header_cells) + " |", "| " + " | ".join(divider) + " |"]
                    for r in body:
                        row_cells = [norm(c) for c in r]
                        # pad/truncate to header length
                        if len(row_cells) < len(header_cells):
                            row_cells += [""] * (len(header_cells) - len(row_cells))
                        table_lines.append("| " + " | ".join(row_cells[: len(header_cells)]) + " |")
                    processed_rows = len(body)
                    skipped_empty_rows = 0

                sheet_md = f"### {sheet.title}\n\n" + "\n".join(table_lines)

                # 処理結果のログ情報（コメントとして追加）
                if empty_line_config.enabled and skipped_empty_rows > 0:
                    log_info = f"\n<!-- 空白行処理: {original_row_count}行 → {processed_rows}行 (空白行{skipped_empty_rows}行をスキップ) -->"
                    sheet_md += log_info
            else:
                # single column or sparse: render each non-empty row as a bullet list
                original_line_count = len(rows)
                processed_lines = 0

                try:
                    lines = []
                    for r in rows:
                        if r and r[0] is not None:
                            # 空白行処理が有効で、空白行の場合はスキップ
                            if empty_line_config.enabled and is_empty_cell(r[0]):
                                continue
                            lines.append(str(r[0]))
                            processed_lines += 1

                    sheet_md = f"### {sheet.title}\n\n" + "\n".join([f"- {l}" for l in lines])

                    # 処理結果のログ情報（単一列の場合）
                    if empty_line_config.enabled and (original_line_count - processed_lines) > 0:
                        skipped_lines = original_line_count - processed_lines
                        log_info = f"\n<!-- 空白行処理: {original_line_count}行 → {processed_lines}行 (空白行{skipped_lines}行をスキップ) -->"
                        sheet_md += log_info

                except Exception as exc:
                    # 空白行処理でエラーが発生した場合、元の処理にフォールバック
                    lines = [str(r[0]) for r in rows if r and r[0] is not None]
                    sheet_md = f"### {sheet.title}\n\n" + "\n".join([f"- {l}" for l in lines])

            sheets_md.append(sheet_md)

        logger.info(f"Excel XLSX/XLSMファイル変換完了: {file_name} ({len(sheets_md)}シート処理)")
        return "\n\n".join(sheets_md)

    # --- DOC ---
    if ext == ".doc":
        # .doc files are also supported by python-docx in newer versions
        try:
            import docx  # type: ignore
        except Exception as exc:
            raise RuntimeError("python-docx is required to convert .doc files") from exc

        try:
            empty_line_config = get_empty_line_config()
            doc = docx.Document(path)  # type: ignore

            try:
                if empty_line_config.enabled:
                    # 空白行処理が有効な場合、段落をフィルタリング
                    paragraphs = []
                    for p in doc.paragraphs:
                        if p.text and not is_empty_cell(p.text.strip()):
                            paragraphs.append(p.text)
                        elif not empty_line_config.remove_consecutive:
                            # 連続する空白行の除去が無効の場合、空の段落も含める
                            paragraphs.append('')
                else:
                    paragraphs = [p.text for p in doc.paragraphs if p.text]

                # 末尾の空白行処理
                if empty_line_config.enabled and empty_line_config.remove_trailing:
                    while paragraphs and is_empty_cell(paragraphs[-1]):
                        paragraphs.pop()

                raw = "\n\n".join(paragraphs)
            except Exception as exc:
                # 空白行処理でエラーが発生した場合、元の処理にフォールバック
                logger.warning(f"DOC空白行処理でエラー、フォールバック実行: {exc}")
                paragraphs = [p.text for p in doc.paragraphs if p.text]
                raw = "\n\n".join(paragraphs)

            logger.info(f"DOCファイル変換完了: {file_name}")
            return _maybe_markitdown_convert(raw)
        except Exception as exc:
            raise RuntimeError(f"Failed to convert .doc file: {exc}") from exc

    # --- XLS ---
    if ext == ".xls":
        try:
            import xlrd  # type: ignore
        except Exception as exc:
            raise RuntimeError("xlrd is required to convert .xls files") from exc

        try:
            workbook = xlrd.open_workbook(path)
            sheets_md: List[str] = []
            empty_line_config = get_empty_line_config()

            for sheet_index in range(workbook.nsheets):
                sheet = workbook.sheet_by_index(sheet_index)
                if sheet.nrows == 0:
                    continue

                rows = []
                for row_idx in range(sheet.nrows):
                    row = [str(sheet.cell_value(row_idx, col_idx)) for col_idx in range(sheet.ncols)]
                    rows.append(row)

                # If there are multiple columns, render as a Markdown table
                if sheet.ncols > 1 and sheet.nrows > 1:
                    original_row_count = len(rows)
                    skipped_empty_rows = 0
                    
                    header = rows[0]
                    body = rows[1:]
                    
                    # 空白行処理が有効な場合、空白行をフィルタリング
                    if empty_line_config.enabled:
                        filtered_body = []
                        for row in body:
                            if is_empty_row(row):
                                skipped_empty_rows += 1
                                continue
                            filtered_body.append(row)
                        body = filtered_body
                        processed_rows = len(body) + 1  # ヘッダー行を含む
                    else:
                        processed_rows = len(body) + 1
                    
                    # divider
                    divider = ["---"] * len(header)
                    table_lines = ["| " + " | ".join(header) + " |", "| " + " | ".join(divider) + " |"]
                    for row in body:
                        # pad/truncate to header length
                        if len(row) < len(header):
                            row += [""] * (len(header) - len(row))
                        table_lines.append("| " + " | ".join(row[:len(header)]) + " |")
                    
                    sheet_md = f"### {sheet.name}\n\n" + "\n".join(table_lines)

                    # 処理結果のログ情報（コメントとして追加）
                    if empty_line_config.enabled and skipped_empty_rows > 0:
                        log_info = f"\n<!-- 空白行処理: {original_row_count}行 → {processed_rows}行 (空白行{skipped_empty_rows}行をスキップ) -->"
                        sheet_md += log_info
                else:
                    # single column or sparse: render as bullet list
                    original_line_count = len(rows)
                    processed_lines = 0

                    lines = []
                    for row in rows:
                        if row and row[0]:
                            # 空白行処理が有効で、空白行の場合はスキップ
                            if empty_line_config.enabled and is_empty_cell(str(row[0])):
                                continue
                            lines.append(str(row[0]))
                            processed_lines += 1

                    sheet_md = f"### {sheet.name}\n\n" + "\n".join([f"- {l}" for l in lines])

                    # 処理結果のログ情報（単一列の場合）
                    if empty_line_config.enabled and (original_line_count - processed_lines) > 0:
                        skipped_lines = original_line_count - processed_lines
                        log_info = f"\n<!-- 空白行処理: {original_line_count}行 → {processed_lines}行 (空白行{skipped_lines}行をスキップ) -->"
                        sheet_md += log_info

                sheets_md.append(sheet_md)

            logger.info(f"Excel XLSファイル変換完了: {file_name} ({len(sheets_md)}シート処理)")
            return "\n\n".join(sheets_md)
        except Exception as exc:
            raise RuntimeError(f"Failed to convert .xls file: {exc}") from exc

    # --- PDF ---
    if ext == ".pdf":
        try:
            import PyPDF2  # type: ignore
        except ImportError:
            try:
                import pypdf  # type: ignore
                PyPDF2 = pypdf
            except ImportError:
                return _maybe_markitdown_convert(path)

        with open(path, "rb") as f:
            reader = PyPDF2.PdfReader(f)
            pages = []
            for page_num, page in enumerate(reader.pages, 1):
                text = page.extract_text()
                if text.strip():
                    pages.append(f"## Page {page_num}\n\n{text}")

        raw = "\n\n".join(pages)
        logger.info(f"PDFファイル変換完了: {file_name} ({len(pages)}ページ)")
        return _maybe_markitdown_convert(raw)

    # --- PPTX ---
    if ext == ".pptx":
        try:
            import pptx  # type: ignore
        except ImportError:
            return _maybe_markitdown_convert(path)

        presentation = pptx.Presentation(path)
        slides_md = []
        for slide_num, slide in enumerate(presentation.slides, 1):
            slide_text = []
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    slide_text.append(shape.text)
            if slide_text:
                slides_md.append(f"### スライド {slide_num}\n\n" + "\n\n".join(slide_text))

        raw = "\n\n".join(slides_md)
        logger.info(f"PowerPointファイル変換完了: {file_name} ({len(slides_md)}スライド)")
        return _maybe_markitdown_convert(raw)

    # --- PPT (Legacy) ---
    if ext == ".ppt":
        # Legacy PowerPoint files - markitdown fallback or error
        file_size = os.path.getsize(path) if os.path.exists(path) else 0
        file_size_mb = file_size / (1024 * 1024)
        
        legacy_info = f"""# PowerPointファイル（レガシー形式）

## ファイル情報
- **ファイル名**: {file_name}
- **ファイルサイズ**: {file_size} bytes
- **最終更新日**: {os.path.getmtime(path) if os.path.exists(path) else 'Unknown'}
- **形式**: Microsoft PowerPoint 97-2003 (.ppt)

## 変換状況
このファイルはレガシーなPowerPoint形式（.ppt）のため、自動的なテキスト抽出は
サポートされていません。ファイルの内容を確認するには、以下の方法をお試しください：

1. Microsoft PowerPointでファイルを開いて.pptx形式で保存し直す
2. LibreOffice Impressでファイルを開いてODP形式で保存する
3. オンライン変換ツールを使用して新しい形式に変換する

## 技術的詳細
- 検出された形式: OLE複合文書形式
- 対応ライブラリ: python-pptx（非対応）、markitdown（非対応）
- 推奨対応: ファイル形式の変換後に再処理
"""
        
        logger.warning(f"PPT（レガシー）形式は直接サポートされていません: {file_name}")
        return legacy_info

    # Fallback to markitdown if available
    logger.info(f"未対応拡張子のため markitdown での変換を試行: {ext}")
    return _maybe_markitdown_convert(path)
